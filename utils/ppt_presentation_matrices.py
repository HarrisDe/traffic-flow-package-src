from sklearn.metrics import mean_absolute_error, median_absolute_error, mean_squared_error, mean_absolute_percentage_error
from sklearn.preprocessing import StandardScaler, MinMaxScaler
import numpy as np
import pandas as pd
import pickle
import os
import glob
from sklearn.preprocessing import StandardScaler, MinMaxScaler
import numpy as np
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import subprocess
from typing import List, Dict, Optional, Union, Tuple
import warnings 



class PPTMatrixGenerator:
    """
    Class to generate PowerPoint presentations with matrices/tables
    based on provided CSV or DataFrame, fully configurable.
    """

    def __init__(self,
                 data: Union[str, pd.DataFrame],
                 columns: List[str],
                 ppt_name: str,
                 group_by: Optional[List[str]] = None,
                 two_tables_per_slide: bool = False,
                 four_tables_per_slide: bool = False,
                 model_column: Optional[str] = None,
                 show_std_in_mae: bool = True,
                 column_renames: Optional[Dict[str, str]] = None):
        self.data = data
        self.columns = columns
        self.ppt_name = ppt_name
        self.group_by = group_by or []
        self.two_tables_per_slide = two_tables_per_slide
        self.four_tables_per_slide = four_tables_per_slide
        self.model_column = model_column
        self.show_std_in_mae = show_std_in_mae
        self.column_renames = column_renames or {}

        if self.two_tables_per_slide and self.four_tables_per_slide:
            raise ValueError("Cannot set both two_tables_per_slide and four_tables_per_slide to True!")

        self._load_data()

        missing_columns = [col for col in self.columns if col not in self.df.columns]
        if missing_columns:
            raise ValueError(f"The following columns are missing from the data: {missing_columns}")

        if (self.two_tables_per_slide or self.four_tables_per_slide) and not self.model_column:
            warnings.warn("Splitting tables per slide requires a 'model_column' to split by.", UserWarning)

        for key in self.column_renames.keys():
            if key not in self.columns:
                warnings.warn(f"'{key}' in column_renames is not in the columns list.", UserWarning)

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def _load_data(self) -> None:
        if isinstance(self.data, str):
            self.df = pd.read_csv(self.data)
        else:
            self.df = self.data.copy()

    def _create_slide(self,
                      base_title: str,
                      model_keys: Optional[List[Union[str, int]]],
                      models: List[pd.DataFrame]) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = base_title

        if len(models) == 1:
            lefts, tops, widths, heights = [Inches(1.0)], [Inches(2.0)], [Inches(11.0)], [Inches(4.5)]
        elif len(models) == 2:
            lefts, tops = [Inches(0.7), Inches(6.8)], [Inches(2.0)] * 2
            widths, heights = [Inches(5.5)] * 2, [Inches(3.5)] * 2
        else:
            lefts = [Inches(0.4), Inches(6.9), Inches(0.4), Inches(6.9)]
            tops = [Inches(1.0), Inches(1.0), Inches(3.0), Inches(3.0)]
            widths = [Inches(5.8)] * 4
            heights = [Inches(1.5)] * 4

        for idx, df_sub in enumerate(models):
            if idx >= len(lefts):
                break

            table_data = self._prepare_table_data(df_sub)
            rows, cols = len(table_data) + 1, len(self.columns)
            table = slide.shapes.add_table(rows, cols, lefts[idx], tops[idx], widths[idx], heights[idx]).table

            if model_keys:
                label = str(model_keys[idx])
                label_box = slide.shapes.add_textbox(lefts[idx], tops[idx] - Inches(0.3), widths[idx], Inches(0.3))
                tf = label_box.text_frame
                tf.text = label
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                run = tf.paragraphs[0].runs[0]
                run.font.name = self.font_name
                run.font.size = Pt(self.cell_font_size)

            for col_idx, col_name in enumerate(self.columns):
                cell = table.cell(0, col_idx)
                cell.text = self.column_renames.get(col_name, col_name)
                self._format_cell(cell, header=True)

            for row_idx, row_data in enumerate(table_data, start=1):
                for col_idx, value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = value
                    self._format_cell(cell, header=False)

    def generate(self,
                 font_name: str = "Arial",
                 header_font_size: Optional[int] = None,
                 cell_font_size: Optional[int] = None,
                 alignment: str = "CENTER",
                 open_after: bool = True) -> None:

        if not header_font_size:
            header_font_size = 16 if not (self.two_tables_per_slide or self.four_tables_per_slide) else 14
            header_font_size = 12 if self.four_tables_per_slide else header_font_size

        if not cell_font_size:
            cell_font_size = 14 if not (self.two_tables_per_slide or self.four_tables_per_slide) else 12
            cell_font_size = 10 if self.four_tables_per_slide else cell_font_size

        valid_alignments = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT}
        if alignment.upper() not in valid_alignments:
            warnings.warn("Invalid alignment specified. Defaulting to CENTER.", UserWarning)
            alignment = "CENTER"

        self.alignment = valid_alignments[alignment.upper()]
        self.font_name = font_name
        self.header_font_size = header_font_size
        self.cell_font_size = cell_font_size

        grouped = self.df.groupby(self.group_by) if self.group_by else [(None, self.df)]

        for keys, group in grouped:
            if isinstance(keys, tuple):
                base_title = ", ".join(f"{col}={val}" for col, val in zip(self.group_by, keys))
            elif keys is not None:
                base_title = f"{self.group_by[0]}={keys}"
            else:
                base_title = "All Results"

            if (self.two_tables_per_slide or self.four_tables_per_slide) and self.model_column:
                model_grouped = group.groupby(self.model_column)
                model_keys = list(model_grouped.groups.keys())
                per_slide = 2 if self.two_tables_per_slide else 4

                for i in range(0, len(model_keys), per_slide):
                    selected_keys = model_keys[i:i + per_slide]
                    models = [model_grouped.get_group(k) for k in selected_keys]
                    self._create_slide(base_title, selected_keys, models)
            else:
                self._create_slide(base_title, None, [group])

        output_dir = os.path.dirname(self.ppt_name)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        self.prs.save(self.ppt_name)

        if open_after:
            try:
                subprocess.call(["open", self.ppt_name])
            except Exception:
                try:
                    subprocess.call(["start", self.ppt_name], shell=True)
                except Exception:
                    pass

    def _prepare_table_data(self, df_subset: pd.DataFrame) -> List[List[str]]:
        table_data: List[List[str]] = []
        for _, row in df_subset.iterrows():
            row_data = []
            for col in self.columns:
                if col == "MAE" and self.show_std_in_mae:
                    formatted = f"{row['MAE']:.2f} ± {row['MAE_std']:.1f}"
                else:
                    val = row[col]
                    formatted = f"{val:.2f}" if isinstance(val, float) else str(val)
                row_data.append(formatted)
            table_data.append(row_data)
        return table_data

    def _format_cell(self, cell: _Cell, header: bool = False) -> None:
        cell.vertical_alignment = PP_ALIGN.CENTER
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = self.alignment
            for run in paragraph.runs:
                run.font.name = self.font_name
                run.font.size = Pt(self.header_font_size if header else self.cell_font_size)

    @staticmethod
    def create_combination_column(df: pd.DataFrame,
                                   columns_to_combine: List[str],
                                   new_column_name: str = "combination_label",
                                   insert_first: bool = False) -> pd.DataFrame:
        def combine_row(row) -> str:
            parts = [f"{col}-{row[col]}" for col in columns_to_combine]
            return "_".join(parts)

        df[new_column_name] = df.apply(combine_row, axis=1)

        if insert_first:
            cols = [new_column_name] + [col for col in df.columns if col != new_column_name]
            df = df[cols]

        return df