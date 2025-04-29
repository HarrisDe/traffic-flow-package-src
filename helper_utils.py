from sklearn.metrics import mean_absolute_error, median_absolute_error, mean_squared_error, mean_absolute_percentage_error
from sklearn.preprocessing import StandardScaler, MinMaxScaler
import numpy as np
import pandas as pd
import pickle
import os
import glob
from sklearn.preprocessing import StandardScaler, MinMaxScaler
import numpy as np
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import subprocess
from typing import Union, List, Optional
import warnings 

class LoggingMixin:
    def __init__(self, disable_logs=False):
        self.disable_logs = disable_logs

    def _log(self, message):
        if not self.disable_logs:
            logging.info(message)
            
            
def get_filtered_X(X_train, lags, spatial_adj):
    """
    Filters the input DataFrame by dropping lag and spatial adjacency features 
    that exceed the specified lag and spatial adjacency limits. To be used so
    that all the lags and spatial adjacencies are calculated only once and then 
    the model is evaluated on the filtered dataset..

    Parameters
    ----------
    X_train : pd.DataFrame
        The input DataFrame containing features including temporal lags and spatial adjacency sensors.
    lags : int
        The number of temporal lag features to retain (e.g., if lags=20, keep relative_diff_lag1 to relative_diff_lag20).
    spatial_adj : int
        The number of upstream/downstream adjacent sensor features to retain (e.g., if spatial_adj=3, 
        keep upstream_sensor_1 to upstream_sensor_3 and downstream_sensor_1 to downstream_sensor_3).

    Returns
    -------
    pd.DataFrame
        A filtered DataFrame with only the specified number of lag and spatial adjacency features retained.
    """

    # Make a copy to avoid modifying the original DataFrame
    X = X_train.copy()
    
    # Define the maximum available number of lag features
    max_lags = 30
    if lags < max_lags:
        # Create list of lag features to drop
        drop_lags = [f'relative_diff_lag{i}' for i in range(lags + 1, max_lags + 1)]
        # Drop them if present
        X = X.drop(columns=[col for col in drop_lags if col in X.columns], errors='ignore')
    
    # Define the maximum number of spatial adjacencies (up/down stream sensors)
    max_adj = 5
    if spatial_adj < max_adj:
        # Drop upstream and downstream sensor features beyond the selected number
        for i in range(spatial_adj + 1, max_adj + 1):
            for direction in ['downstream', 'upstream']:
                col = f'{direction}_sensor_{i}'
                if col in X.columns:
                    X = X.drop(columns=col)

    return X


def modify_gman(df_gman_test):
    """
    Prepares GMAN test results for analysis by:
    - Filtering for the final prediction timestep
    - Calculating prediction horizon based on timestamps
    - Reshaping the data to long format for per-sensor predictions

    Assumes:
    - Each prediction sample includes multiple timesteps
    - Prediction columns end with '_pred'

    Args:
        df_gman_test (pd.DataFrame): Wide-format GMAN results DataFrame. Must include
                                     'timestamp' or 'date', 'timestep', 'sample_nr', and prediction columns.

    Returns:
        pd.DataFrame: Long-format DataFrame with columns:
                      - 'sensor_id': sensor identifier
                      - 'gman_prediction': predicted value
                      - 'target_date': date of the forecasted value
                      - 'prediction_date': date when the forecast was made
    """
    # --- Handle timestamp column ---
    if 'timestamp' in df_gman_test.columns:
        date_col = 'timestamp'
    elif 'date' in df_gman_test.columns:
        date_col = 'date'
    else:
        raise ValueError("Expected 'timestamp' or 'date' column in the input DataFrame.")

    df_gman_test[date_col] = pd.to_datetime(df_gman_test[date_col])

    # --- Compute prediction horizon (based on first prediction sample) ---
    first_sample = df_gman_test[df_gman_test['sample_nr'] == 1]
    horizon = first_sample[date_col].max() - first_sample[date_col].min() + pd.Timedelta(minutes=1)

    # --- Keep only the last timestep's predictions ---
    max_timestep = df_gman_test['timestep'].max()
    df_latest = df_gman_test[df_gman_test['timestep'] == max_timestep]

    # --- Extract relevant columns ---
    pred_cols = [col for col in df_latest.columns if col.endswith('_pred')]
    df_latest = df_latest[[date_col] + pred_cols]

    # --- Rename sensor columns (strip '_pred') ---
    df_latest.columns = [date_col] + [col[:-5] for col in pred_cols]

    # --- Reshape to long format ---
    df_long = df_latest.melt(
        id_vars=date_col,
        var_name='sensor_id',
        value_name='gman_prediction'
    )

    # --- Add timestamps for target and prediction ---
    df_long.rename(columns={date_col: 'gman_target_date'}, inplace=True)
    df_long['gman_prediction_date'] = df_long['gman_target_date'] - horizon

    return df_long


def load_gman_results(p, q, directory="saved_gman_results", pattern_template="gman_results_P{p}_Q{q}*.parquet"):
    """
    Load previously saved GMAN results from Parquet files that match a given filename pattern.
    If there are multiple matching files, the first one is selected.

    Args:
        p (int): Number of history steps (including the current timestep).
        q (int): Prediction horizon (number of future steps to forecast).
        directory (str): Directory where results are stored.
        pattern_template (str): Pattern template to match filenames. Should include '{p}' and '{q}' placeholders.

    Returns:
        DataFrame or None: Loaded results DataFrame, or None if no matching files are found.
    """
    # Create full file search pattern using provided template
    file_pattern = os.path.join(directory, pattern_template.format(p=p, q=q))
    
    # Search for matching files
    matching_files = glob.glob(file_pattern)
    if matching_files:
        print(f"Found files: {matching_files}")
        return pd.read_parquet(matching_files[0])  # Return the first match
    else:
        print(f"No files found matching {file_pattern}")
        return None

def normalize_data(X_train=None, X_test=None, use_minmax_norm=False, use_full_data=False):
    """
    Normalizes training and/or testing data using StandardScaler or MinMaxScaler.

    Parameters:
    - X_train (array-like, optional): Training dataset. If None, only `X_test` will be normalized.
    - X_test (array-like, optional): Testing dataset. If None, only `X_train` will be normalized.
    - use_minmax_norm (bool): If True, uses MinMaxScaler; otherwise, uses StandardScaler.
    - use_full_data (bool): If True, normalizes using both training and testing data combined.
        WARNING: Using this option introduces data leakage because the test data influences 
        the scaling applied to the training data. This approach may lead to overly optimistic 
        performance metrics and is not recommended for real-world scenarios where the test 
        set must remain unseen until final evaluation.

    Returns:
    - X_train_normalized (array-like, optional): Normalized training dataset (if `X_train` is provided).
    - X_test_normalized (array-like, optional): Normalized testing dataset (if `X_test` is provided).

    Notes:
    - Default behavior (`use_full_data=False`) ensures that scaling is based solely on the training data, 
      which is a best practice to avoid data leakage.
    - Use `use_full_data=True` only in controlled experiments where you need consistent scaling across 
      both training and test sets and are aware of the potential risks.
    - At least one of `X_train` or `X_test` must be provided.

    Example Usage:
    - Normalize with StandardScaler (default):
        X_train_normalized, X_test_normalized = normalize_data(X_train, X_test)
    - Normalize only `X_train`:
        X_train_normalized, _ = normalize_data(X_train=X_train)
    - Normalize with MinMaxScaler using both datasets:
        X_train_normalized, X_test_normalized = normalize_data(X_train, X_test, use_minmax_norm=True, use_full_data=True)
    """
    if X_train is None and X_test is None:
        raise ValueError("At least one of X_train or X_test must be provided.")

    scaler = MinMaxScaler() if use_minmax_norm else StandardScaler()

    if use_full_data:
        if X_train is not None and X_test is not None:
            # Concatenate training and test data for joint scaling
            full_data = np.concatenate([X_train, X_test], axis=0)
            full_data_normalized = scaler.fit_transform(full_data)
            # Split back into training and test sets
            X_train_normalized = full_data_normalized[:X_train.shape[0], :]
            X_test_normalized = full_data_normalized[X_train.shape[0]:, :]
        elif X_train is not None:
            X_train_normalized = scaler.fit_transform(X_train)
            X_test_normalized = None
        elif X_test is not None:
            X_test_normalized = scaler.fit_transform(X_test)
            X_train_normalized = None
    else:
        if X_train is not None:
            X_train_normalized = scaler.fit_transform(X_train)
            if X_test is not None:
                X_test_normalized = scaler.transform(X_test)
            else:
                X_test_normalized = None
        elif X_test is not None:
            # Scale only test data if training data is not provided
            X_test_normalized = scaler.fit_transform(X_test)
            X_train_normalized = None
    

    return X_train_normalized, X_test_normalized


def load_and_evaluate_a_model(model_name, X_train, X_test, y_train, y_test, horizon, return_csv=True):
    """
    Load the best model by name, make predictions on a new dataset, and calculate prediction errors
    along with naive predictions.

    Parameters:
    - model_name (str): The name of the model to load (e.g., 'XGBoost', 'Random_Forest', 'Neural_Network').
    - new_X (pd.DataFrame): New input features for prediction.
    - new_y (pd.Series): True target values for the new dataset.
    - use_normalized (bool): If True, normalize the input features for ANN.

    Returns:
    - dict: A dictionary containing model MAE and naive MAE for comparison.
    """
    # Load the model
    model_path = f'../models/{model_name}'
    if 'neural' in model_name.lower():
        from keras.models import load_model
        best_model = load_model(f'{model_path}.h5')
        print(f"{model_name} model loaded from {model_path}.h5")
        print("Normalizing new dataset for ANN.")
        X_train_normalized, X_test_normalized = normalize_data(
            X_train, X_test, use_minmax_norm=False, use_full_data=False)
        y_pred = best_model.predict(X_test_normalized)

    else:
        with open(f'{model_path}.pkl', 'rb') as f:
            best_model = pickle.load(f)
        print(f"{model_name} model loaded from {model_path}.pkl")
        y_pred = best_model.predict(X_test)

    test_mae = abs(y_test - y_pred).mean()

    # Naive model: Predict the last value (last observation before the prediction)
    naive_predictions = np.abs(y_test)  # (naive(same as last) - actual)
    naive_mae = np.mean(naive_predictions)

    print(f'Test MAE (horizon: {horizon}min):{test_mae:.4f}')
    print(f'Test naive MAE (horizon: {horizon}min:{naive_mae:.4f}')

    if return_csv:
        df = {'incremental_id': X_test['incremental_id'],
              f'y_test_{horizon}': y_test, f'y_pred_{horizon}': y_pred}
        df = pd.DataFrame(df)
        df.to_csv(f'results_horizon_{horizon}_min.csv')


def optimize_dtypes(df):
    """
    Optimize the data types of numeric columns in a DataFrame to reduce memory usage.

    This function identifies numeric columns (`int` and `float`) in the DataFrame and
    downcasts their data types to the smallest possible type (`float32`, `int32`, etc.)
    without losing precision. This can help reduce memory usage for large DataFrames.

    Parameters:
    - df (pd.DataFrame): The input DataFrame to optimize.

    Returns:
    - pd.DataFrame: The DataFrame with optimized data types for numeric columns.
    """

    # Iterate over numeric columns (int and float types)
    for col in df.select_dtypes(include=['int', 'float']).columns:
        # Check and downcast float columns
        if df[col].dtype == 'float64':
            df[col] = pd.to_numeric(df[col], downcast='float')
        # Check and downcast integer columns
        elif df[col].dtype == 'int64':
            df[col] = pd.to_numeric(df[col], downcast='integer')

    return df



class PPTMatrixGenerator:
    """
    Class to generate PowerPoint presentations with matrices/tables
    based on provided CSV or DataFrame, fully configurable.
    """

    def __init__(self,
                 data,
                 columns,
                 ppt_name,
                 group_by=None,
                 two_tables_per_slide=False,
                 four_tables_per_slide=False,
                 model_column=None,
                 show_std_in_mae=True):
        """
        Initialize the generator.

        Args:
            data (str or pd.DataFrame): Path to CSV file or loaded DataFrame.
            columns (list): List of columns to include in tables.
            ppt_name (str): Name of the output PowerPoint file.
            group_by (list, optional): Columns to group by. Defaults to None.
            two_tables_per_slide (bool, optional): Whether to split two tables per slide.
            four_tables_per_slide (bool, optional): Whether to split four tables per slide (2x2 grid).
            model_column (str, optional): Column to split inside group for multiple tables per slide.
            show_std_in_mae (bool, optional): Whether to display MAE ± MAE_std format. Defaults to True.
        """
        self.data = data
        self.columns = columns
        self.ppt_name = ppt_name
        self.group_by = group_by or []
        self.two_tables_per_slide = two_tables_per_slide
        self.four_tables_per_slide = four_tables_per_slide
        self.model_column = model_column
        self.show_std_in_mae = show_std_in_mae

        # Basic validation
        if self.two_tables_per_slide and self.four_tables_per_slide:
            raise ValueError("Cannot set both two_tables_per_slide and four_tables_per_slide to True!")

        self._load_data()

        # Check columns exist
        missing_columns = [col for col in self.columns if col not in self.df.columns]
        if missing_columns:
            raise ValueError(f"The following columns are missing from the data: {missing_columns}")

        # Check if splitting without model_column
        if (self.two_tables_per_slide or self.four_tables_per_slide) and not self.model_column:
            warnings.warn("Splitting tables per slide requires a 'model_column' to split by.", UserWarning)

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def _load_data(self):
        """Load the data if a path is provided."""
        if isinstance(self.data, str):
            self.df = pd.read_csv(self.data)
        else:
            self.df = self.data.copy()
            
    def _create_slide(self, base_title, model_keys, models):
        """Create slide with 1, 2, or 4 matrices with labels."""

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = base_title

        if len(models) == 1:
            lefts = [Inches(1.0)]
            tops = [Inches(2.0)]
            widths = [Inches(11.0)]
            heights = [Inches(4.5)]
        elif len(models) == 2:
            lefts = [Inches(0.7), Inches(6.8)]
            tops = [Inches(2.0), Inches(2.0)]
            widths = [Inches(5.5)] * 2
            heights = [Inches(3.5)] * 2
        else:
            lefts = [Inches(0.4), Inches(6.9), Inches(0.4), Inches(6.9)]
            tops = [Inches(1.0), Inches(1.0), Inches(3.0), Inches(3.0)]
            widths = [Inches(5.8)] * 4
            heights = [Inches(1.5)] * 4

        for idx, df_sub in enumerate(models):
            if idx >= len(lefts):
                break

            table_data = self._prepare_table_data(df_sub)
            rows, cols = len(table_data) + 1, len(self.columns)

            # Add the table
            table = slide.shapes.add_table(rows, cols, lefts[idx], tops[idx], widths[idx], heights[idx]).table

            # Add label above each table
            if model_keys:
                label = str(model_keys[idx])
                label_box = slide.shapes.add_textbox(
                    lefts[idx],
                    tops[idx] - Inches(0.3),
                    widths[idx],
                    Inches(0.3)
                )
                tf = label_box.text_frame
                tf.text = label
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                run = tf.paragraphs[0].runs[0]
                run.font.name = self.font_name
                run.font.size = Pt(self.cell_font_size)

            # Table header
            for col_idx, col_name in enumerate(self.columns):
                cell = table.cell(0, col_idx)
                cell.text = col_name
                self._format_cell(cell, header=True)

            # Table data
            for row_idx, row_data in enumerate(table_data, start=1):
                for col_idx, value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = value
                    self._format_cell(cell, header=False)

    def generate(self,
                 font_name="Arial",
                 header_font_size=None,
                 cell_font_size=None,
                 alignment="CENTER",
                 open_after=True):
        """
        Generate the PowerPoint presentation.

        Args:
            font_name (str, optional): Font name for text.
            header_font_size (int, optional): Header font size.
            cell_font_size (int, optional): Cell font size.
            alignment (str, optional): LEFT, CENTER, or RIGHT.
            open_after (bool, optional): Automatically open the PPT after creation.
        """

        # Auto-set font sizes
        if not header_font_size:
            header_font_size = 16 if not (self.two_tables_per_slide or self.four_tables_per_slide) else 14
            header_font_size = 12 if self.four_tables_per_slide else header_font_size

        if not cell_font_size:
            cell_font_size = 14 if not (self.two_tables_per_slide or self.four_tables_per_slide) else 12
            cell_font_size = 10 if self.four_tables_per_slide else cell_font_size

        # Validate alignment
        valid_alignments = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT}
        if alignment.upper() not in valid_alignments:
            warnings.warn("Invalid alignment specified. Defaulting to CENTER.", UserWarning)
            alignment = "CENTER"

        self.alignment = valid_alignments[alignment.upper()]
        self.font_name = font_name
        self.header_font_size = header_font_size
        self.cell_font_size = cell_font_size

        # Group or not
        if self.group_by:
            grouped = self.df.groupby(self.group_by)
        else:
            grouped = [(None, self.df)]

        for keys, group in grouped:
            if isinstance(keys, tuple):
                base_title = ", ".join(f"{col}={val}" for col, val in zip(self.group_by, keys))
            elif keys is not None:
                base_title = f"{self.group_by[0]}={keys}"
            else:
                base_title = "All Results"

            if (self.two_tables_per_slide or self.four_tables_per_slide) and self.model_column:
                model_grouped = group.groupby(self.model_column)
                model_keys = list(model_grouped.groups.keys())
                per_slide = 2 if self.two_tables_per_slide else 4

                for i in range(0, len(model_keys), per_slide):
                    selected_keys = model_keys[i:i+per_slide]
                    models = [model_grouped.get_group(k) for k in selected_keys]
                    self._create_slide(base_title, selected_keys, models)
            else:
                self._create_slide(base_title, None, [group])

        self.prs.save(self.ppt_name)
        if open_after:
            try:
                subprocess.call(["open", self.ppt_name])  # macOS/Linux
            except Exception:
                try:
                    subprocess.call(["start", self.ppt_name], shell=True)  # Windows
                except Exception:
                    pass

    def _prepare_table_data(self, df_subset):
        """Prepare formatted rows for the table."""
        table_data = []
        for _, row in df_subset.iterrows():
            row_data = []
            for col in self.columns:
                if col == "MAE" and self.show_std_in_mae:
                    formatted = f"{row['MAE']:.2f} ± {row['MAE_std']:.1f}"
                else:
                    val = row[col]
                    formatted = f"{val:.2f}" if isinstance(val, float) else str(val)
                row_data.append(formatted)
            table_data.append(row_data)
        return table_data

    

    def _format_cell(self, cell, header=False):
        """Format a cell with font and alignment."""
        cell.vertical_alignment = PP_ALIGN.CENTER
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = self.alignment
            for run in paragraph.runs:
                run.font.name = self.font_name
                run.font.size = Pt(self.header_font_size if header else self.cell_font_size)

    @staticmethod
    def create_combination_column(df, columns_to_combine, new_column_name="combination_label", insert_first=False):
        """
        Create a new combined label column based on multiple columns.

        Args:
            df (pd.DataFrame): Input dataframe.
            columns_to_combine (list of str): Columns to combine.
            new_column_name (str): Name for the new column.
            insert_first (bool): Whether to insert the new column as first column.

        Returns:
            pd.DataFrame
        """
        def combine_row(row):
            parts = [f"{col}-{row[col]}" for col in columns_to_combine]
            return "_".join(parts)

        df[new_column_name] = df.apply(combine_row, axis=1)

        if insert_first:
            cols = [new_column_name] + [col for col in df.columns if col != new_column_name]
            df = df[cols]

        return df